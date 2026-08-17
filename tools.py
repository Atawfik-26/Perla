import os
import re
import uuid

from docx import Document
from docx.shared import Pt

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.utils import simpleSplit

from pptx import Presentation

try:
    from openpyxl import Workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False


# =========================================================
# CONFIG
# =========================================================

FILES_DIR = "files"
os.makedirs(FILES_DIR, exist_ok=True)


# =========================================================
# HELPERS
# =========================================================

def _clean_filename(text, default="perla_file"):
    text = (text or default).strip()

    text = re.sub(
        r"[^\w\u0600-\u06FF\- ]",
        "",
        text
    )

    text = text.replace(" ", "_")

    return text[:40] or default


def _find_arabic_font():
    """
    يحاول العثور على خط Unicode موجود على ويندوز
    يدعم العربية.
    """

    possible_fonts = [
        r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\Arial.ttf",
        r"C:\Windows\Fonts\tahoma.ttf",
        r"C:\Windows\Fonts\Tahoma.ttf",
        r"C:\Windows\Fonts\segoeui.ttf",
        r"C:\Windows\Fonts\segoeui.ttf",
    ]

    for font_path in possible_fonts:
        if os.path.exists(font_path):
            return font_path

    return None


PDF_FONT_NAME = "Helvetica"

arabic_font_path = _find_arabic_font()

if arabic_font_path:
    try:
        pdfmetrics.registerFont(
            TTFont(
                "PerlaArabic",
                arabic_font_path
            )
        )

        PDF_FONT_NAME = "PerlaArabic"

    except Exception as error:
        print(
            f"[PERLA PDF] فشل تحميل الخط العربي: {repr(error)}"
        )


# =========================================================
# WORD / DOCX
# =========================================================

def create_docx(content, title=None):
    """
    إنشاء ملف Word.
    """

    content = content or ""

    doc = Document()

    # العنوان
    if title:
        heading = doc.add_heading(
            title,
            level=1
        )

        for run in heading.runs:
            run.font.size = Pt(18)

    # المحتوى
    for paragraph in content.split("\n"):

        paragraph = paragraph.strip()

        if paragraph:

            p = doc.add_paragraph(
                paragraph
            )

            for run in p.runs:
                run.font.size = Pt(12)

        else:
            doc.add_paragraph("")

    filename = (
        f"{_clean_filename(title)}_"
        f"{uuid.uuid4().hex[:6]}.docx"
    )

    path = os.path.join(
        FILES_DIR,
        filename
    )

    doc.save(path)

    return path, filename


# =========================================================
# PDF
# =========================================================

def create_pdf(content, title=None):
    """
    إنشاء ملف PDF.

    يستخدم خط Unicode إذا وجد خط مناسب
    على الجهاز.
    """

    content = content or ""

    filename = (
        f"{_clean_filename(title, 'perla_pdf')}_"
        f"{uuid.uuid4().hex[:6]}.pdf"
    )

    path = os.path.join(
        FILES_DIR,
        filename
    )

    c = canvas.Canvas(
        path,
        pagesize=A4
    )

    width, height = A4

    margin = 2 * cm

    y = height - margin

    font_size = 12

    font_name = PDF_FONT_NAME

    c.setFont(
        font_name,
        font_size
    )

    max_width = width - (2 * margin)

    # -----------------------------------------------------
    # TITLE
    # -----------------------------------------------------

    if title:

        c.setFont(
            font_name,
            16
        )

        title_lines = simpleSplit(
            title,
            font_name,
            16,
            max_width
        )

        for line in title_lines:

            if y < margin:
                c.showPage()

                c.setFont(
                    font_name,
                    16
                )

                y = height - margin

            c.drawString(
                margin,
                y,
                line
            )

            y -= 0.8 * cm

        y -= 0.5 * cm

        c.setFont(
            font_name,
            font_size
        )

    # -----------------------------------------------------
    # CONTENT
    # -----------------------------------------------------

    for raw_line in content.split("\n"):

        raw_line = raw_line.strip()

        # سطر فاضي
        if not raw_line:

            y -= 0.4 * cm

            if y < margin:

                c.showPage()

                c.setFont(
                    font_name,
                    font_size
                )

                y = height - margin

            continue

        lines = simpleSplit(
            raw_line,
            font_name,
            font_size,
            max_width
        )

        if not lines:
            lines = [""]

        for line in lines:

            if y < margin:

                c.showPage()

                c.setFont(
                    font_name,
                    font_size
                )

                y = height - margin

            c.drawString(
                margin,
                y,
                line
            )

            y -= 0.6 * cm

    c.save()

    return path, filename


# =========================================================
# POWERPOINT
# =========================================================

def create_pptx(content, title=None):
    """
    إنشاء عرض PowerPoint.
    """

    content = content or ""

    prs = Presentation()

    # -----------------------------------------------------
    # تقسيم المحتوى إلى أجزاء
    # -----------------------------------------------------

    parts = [
        p.strip()
        for p in content.split("\n\n")
        if p.strip()
    ]

    slides_data = [
        {
            "title": f"نقطة {i + 1}",
            "content": part
        }

        for i, part in enumerate(parts)
    ]

    if not slides_data:

        slides_data = [
            {
                "title": title or "بيرلا",
                "content": content
            }
        ]

    # -----------------------------------------------------
    # TITLE SLIDE
    # -----------------------------------------------------

    title_slide_layout = prs.slide_layouts[0]

    title_slide = prs.slides.add_slide(
        title_slide_layout
    )

    if title_slide.shapes.title:

        title_slide.shapes.title.text = (
            title or "عرض تقديمي"
        )

    # -----------------------------------------------------
    # CONTENT SLIDES
    # -----------------------------------------------------

    content_slide_layout = prs.slide_layouts[1]

    for item in slides_data:

        slide = prs.slides.add_slide(
            content_slide_layout
        )

        if slide.shapes.title:

            slide.shapes.title.text = (
                item["title"]
            )

        if len(slide.placeholders) > 1:

            text_frame = (
                slide.placeholders[1]
                .text_frame
            )

            text_frame.clear()

            text_frame.text = (
                item["content"]
            )

    # -----------------------------------------------------
    # SAVE
    # -----------------------------------------------------

    filename = (
        f"{_clean_filename(title, 'perla_pptx')}_"
        f"{uuid.uuid4().hex[:6]}.pptx"
    )

    path = os.path.join(
        FILES_DIR,
        filename
    )

    prs.save(path)

    return path, filename


# =========================================================
# EXCEL
# =========================================================

def create_xlsx(content, title=None):
    """
    إنشاء ملف Excel بسيط.

    كل سطر في المحتوى يصبح صفًا.
    لو السطر يحتوي على | يتم تقسيمه إلى أعمدة.
    """

    if not OPENPYXL_AVAILABLE:

        raise RuntimeError(
            "مكتبة openpyxl مش مثبتة. "
            "ثبتها بالأمر: pip install openpyxl"
        )

    content = content or ""

    workbook = Workbook()

    worksheet = workbook.active

    worksheet.title = "بيرلا"

    # -----------------------------------------------------
    # TITLE
    # -----------------------------------------------------

    if title:

        worksheet.append([
            title
        ])

    # -----------------------------------------------------
    # CONTENT
    # -----------------------------------------------------

    lines = content.splitlines()

    for line in lines:

        line = line.strip()

        if not line:
            continue

        # دعم الجداول البسيطة:
        # اسم | السن | الوظيفة

        if "|" in line:

            columns = [
                cell.strip()
                for cell in line.split("|")
            ]

            worksheet.append(
                columns
            )

        else:

            worksheet.append([
                line
            ])

    # -----------------------------------------------------
    # COLUMN WIDTH
    # -----------------------------------------------------

    for column_cells in worksheet.columns:

        max_length = 0

        column_letter = (
            column_cells[0].column_letter
        )

        for cell in column_cells:

            value = str(
                cell.value or ""
            )

            if len(value) > max_length:
                max_length = len(value)

        worksheet.column_dimensions[
            column_letter
        ].width = min(
            max(max_length + 2, 12),
            50
        )

    # -----------------------------------------------------
    # SAVE
    # -----------------------------------------------------

    filename = (
        f"{_clean_filename(title, 'perla_excel')}_"
        f"{uuid.uuid4().hex[:6]}.xlsx"
    )

    path = os.path.join(
        FILES_DIR,
        filename
    )

    workbook.save(path)

    return path, filename
